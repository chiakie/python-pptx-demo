from pptx import Presentation
from pptx.util import Inches, Pt, Cm 
# from pptx.enum.shapes import MSO_SHAPE

import os
from pathlib import Path
# import requests

import pandas as pd

TITLE_SLIDE_LAYOUT = 4
DETAIL_SLIDE_LAYOUT = 6
GROUP_SLIDE_LAYOUT = 7
PERSONAL_SLIDE_LAYOUT = 8

def insetImage(slide, id):
    img_dir = Path(f"images/{id}/")
    left = 1
    top = 3.3
    width = 8
    for file_path in img_dir.rglob('*'):  # 或 rglob('*.txt') 篩選特定副檔名
        if file_path.is_file():
            try:
                slide.shapes.add_picture(str(file_path), Cm(left), Cm(top), Cm(width))
            except (ValueError, TypeError):
                print(f"{id} has ValueError!")
            except Exception as e:
                # 处理所有其他异常
                print(f"发生了一个未知的错误: {e}")
            finally:
                left=left+0.5
                top=top+0.5
                continue

def createTitlePage(prs, id, prize, name, title, unit):
    title_slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    shapes[0].text_frame.text = prize
    shapes[1].text_frame.text = name
    shapes[2].text_frame.text = title
    shapes[3].text_frame.text = unit

    insetImage(slide, id)

def createDetailPage(prs, id, prize, nameTitle, motto, achievement):
    detail_slide_layout = prs.slide_layouts[DETAIL_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(detail_slide_layout)

    # for shape in slide.shapes:
    #     # print(f"  Shape ID: {shape.shape_id}, Type: {shape.shape_type}")
    #     sp = shape._sp
    #     sp.getparent().remove(sp)

    # left = top = width = height = Pt(1)
    # tb = slide.shapes.add_textbox(left, top, width, height)
    # tbtf = tb.text_frame

    shapes = slide.shapes
    shapes[0].text_frame.text = prize
    shapes[4].text_frame.text = "★座右銘："+motto
    shapes[3].text_frame.text = nameTitle
    shapes[2].text_frame.text = "★重點事蹟："

    btf1 = shapes[1].text_frame
    achieves = achievement.splitlines()
    for idx, ach in enumerate(achieves):
        if idx == 0:
            btf1.text = ach
        else:
            p = btf1.add_paragraph()
            p.text = ach
            p.level = 0
    
    insetImage(slide, id)

def createGroupPage(prs, id, prize, subprize, unit, achievement):
    group_slide_layout = prs.slide_layouts[GROUP_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(group_slide_layout)

    shapes = slide.shapes
    shapes[0].text_frame.text = prize
    shapes[1].text_frame.text = subprize
    shapes[2].text_frame.text = unit
    shapes[4].text_frame.text = "★優良事蹟："
    shapes[3].text_frame.text = achievement

    insetImage(slide, id)

def createPersonalPage(prs, id, prize, subprize, unit, name, motto, achievement):
    personal_slide_layout = prs.slide_layouts[PERSONAL_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(personal_slide_layout)

    shapes = slide.shapes
    shapes[0].text_frame.text = prize
    shapes[1].text_frame.text = subprize
    shapes[2].text_frame.text = name
    shapes[4].text_frame.text = unit
    shapes[5].text_frame.text = "★座右銘："+motto
    shapes[3].text_frame.text = "★重點事蹟："
    shapes[6].text_frame.text = achievement

    insetImage(slide, id)

# def downloadImage(urlStr, id):
#     save_dir = f'images/{id}'
#     os.makedirs(save_dir, exist_ok=True)  # 若資料夾不存在就建立

#     urls = urlStr.split(";")
#     for idx, url in enumerate(urls):
#         print(url)
#         response = requests.get(url)
#         if response.status_code == 200:
#             filename = f'{id}_{idx+1}'
#             file_path = os.path.join(save_dir, filename)
#             with open(file_path, 'wb') as f:
#                 f.write(response.content)
#         else:
#             print(f"下載失敗，狀態碼: {response.status_code}")

if __name__ == '__main__':
    prs = Presentation("demo.pptx")

    # for i in range(133):
    #     os.makedirs(f"images/{i+1}", exist_ok=True)  # 若資料夾不存在就建立

    # 讀取Sheet1
    df1 = pd.read_excel("datasource.xlsx", sheet_name='Sheet1', engine='openpyxl')
    for index, row in df1.iterrows():
        id = row[0]
        prize = row[5]
        name = " ".join(row[8])
        imgUrlStrs = row[9]
        motto = row[10]
        achievement = row[11]
        title = " ".join(row[4])
        unit = row[29]

        createTitlePage(prs, id, prize, name, title, unit)
        createDetailPage(prs, id, prize, name+" "+title, motto, achievement)

    # 讀取Sheet2
    df2 = pd.read_excel("datasource.xlsx", sheet_name='Sheet2', engine='openpyxl')
    for index, row in df2.iterrows():
        id = row[0]
        prize = row[5]
        subprize = row[6]
        awardeeType = row[7]
        unit = ""
        motto = ""
        achievement = ""
        if awardeeType == "單位":
            unit = row[12]
            achievement = row[16]
            createGroupPage(prs, id, prize, subprize, unit, achievement)
        else:
            unit = row[29]
            achievement = row[11]
            name = row[8]
            motto = row[10]
            createPersonalPage(prs, id, prize, subprize, unit, name, motto, achievement)     

    prs.save('output.pptx')